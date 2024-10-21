import threading

from flask import Flask, request, jsonify
import os
from io import BytesIO
import numpy as np
import faiss
from bs4 import BeautifulSoup
import requests
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import fitz  # PyMuPDF for PDF processing
from dotenv import load_dotenv
import logging
import streamlit as st
import requests
from io import BytesIO
from flask_cors import CORS

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize the Flask app
api = Flask(__name__)
CORS(api)  # Enable CORS for all routes
load_dotenv()


# Global variables
documents = []
faiss_index = None

# Load environment variables
os.environ["AZURE_OPENAI_API_KEY"] = os.getenv("AZURE_OPENAI_API_KEY_GPT4")
# Add additional environment variables as needed

# Initialize language model and embeddings
llm = AzureChatOpenAI(
    openai_api_version=os.getenv("OPENAI_API_VERSION"),
    azure_deployment=os.getenv("AZURE_DEPLOYMENT_GPT4"),
    model_name=os.getenv("MODEL_NAME"),
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT_GPT4"),
    temperature=0.0,
)

os.environ["AZURE_OPENAI_API_KEY"] = os.getenv("AZURE_OPENAI_API_KEY_EMBEDDING")
os.environ["AZURE_OPENAI_ENDPOINT"] = os.getenv("AZURE_OPENAI_ENDPOINT_EMBEDDING")
os.environ["AZURE_OPENAI_API_VERSION"] = os.getenv("AZURE_OPENAI_API_VERSION_EMBEDDING")

embeddings = AzureOpenAIEmbeddings(
    azure_deployment=os.getenv("AZURE_EMBEDDING_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION_EMBEDDING")
)

# Function to run Flask app
def run_flask():
    api.run(port=5000)

# Start Flask in a separate thread
threading.Thread(target=run_flask, daemon=True).start()

# Streamlit UI
st.title("DocsTalk: Talk to your documents")

# File upload section
st.header("Upload Documents")
uploaded_files = st.file_uploader("Choose files", accept_multiple_files=True, type=['txt', 'docx', 'xlsx', 'pptx', 'pdf'])

if uploaded_files:
    for file in uploaded_files:
        file_content = BytesIO(file.read())
        files = {'file': (file.name, file_content, file.type)}
        response = requests.post('http://127.0.0.1:5000/upload', files=files)
        if response.status_code == 200:
            st.success(f"File {file.name} uploaded and processed successfully!")
        else:
            error_message = response.json().get('error', 'Error uploading file.')
            st.error(f"Error uploading file {file.name}: {error_message}")

# URL input section
st.header("Add URLs")
url_input = st.text_input("Enter URLs (comma-separated)")
if st.button("Process URLs"):
    urls = [url.strip() for url in url_input.split(',') if url.strip()]
    if urls:
        response = requests.post('http://127.0.0.1:5000/process_urls', json={'urls': urls})
        if response.status_code == 200:
            st.success("All URLs processed successfully!")
        else:
            error_message = response.json().get('error', 'Error processing URLs.')
            st.error(error_message)
    else:
        st.warning("Please enter at least one URL.")

# Q&A section
st.header("Ask a Question")
question = st.text_input("Enter your question")
if st.button("Get Answer"):
    if question:
        response = requests.post('http://127.0.0.1:5000/ask', json={'question': question})
        if response.status_code == 200:
            answer = response.json()['answer']
            st.write("Answer:", answer)
        else:
            error_message = response.json().get('error', 'Error getting answer.')
            st.error(error_message)
    else:
        st.warning("Please enter a question.")

# Flask API endpoints

@api.route('/upload', methods=['POST'])
def upload_file():
    global documents, faiss_index
    logger.info("Received file upload request.")
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    try:
        content = file.read()
        file_extension = os.path.splitext(file.filename)[1].lower()
        logger.info(f"Processing file: {file.filename} with extension: {file_extension}")

        # Process different file types
        new_docs = []
        if file_extension == '.txt':
            new_docs = process_text(content.decode('utf-8'), file.filename)
        elif file_extension == '.docx':
            doc = DocxDocument(BytesIO(content))
            text = "\n".join([para.text for para in doc.paragraphs if para.text])
            new_docs = process_text(text, file.filename)
        elif file_extension == '.xlsx':
            df = pd.read_excel(BytesIO(content))
            text = df.to_string(index=False)
            new_docs = process_text(text, file.filename)
        elif file_extension == '.pptx':
            prs = Presentation(BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            new_docs = process_text(text, file.filename)
        elif file_extension == '.pdf':
            new_docs = process_pdf(content, file.filename)
        else:
            return jsonify({"error": f"Unsupported file type: {file_extension}"}), 400

        documents.extend(new_docs)
        update_faiss_index()
        logger.info(f"Successfully processed {file.filename}.")
        return jsonify({"message": "File processed successfully"}), 200

    except Exception as e:
        logger.exception("An error occurred while processing the file.")
        return jsonify({"error": "An error occurred while processing the file."}), 500

@api.route('/process_urls', methods=['POST'])
def process_urls():
    global documents, faiss_index
    urls = request.json.get('urls', [])
    errors = []

    for url in urls:
        if not url.startswith(('http://', 'https://')):
            errors.append(f"Invalid URL '{url}'")
            continue

        try:
            response = requests.get(url)
            if response.status_code == 200:
                text_content = response.text
                new_docs = process_text(text_content, url)
                documents.extend(new_docs)
            else:
                errors.append(f"Error fetching URL '{url}': Status code {response.status_code}.")
        except Exception as e:
            errors.append(f"Error processing URL '{url}': {str(e)}")

    update_faiss_index()

    if errors:
        return jsonify({"error": errors}), 400

    return jsonify({"message": "URLs processed successfully."}), 200

def update_faiss_index():
    global faiss_index
    if documents:
        embeddings_matrix = np.array(embeddings.embed_documents([doc.page_content for doc in documents]))
        embedding_dimension = embeddings_matrix.shape[1]
        faiss_index = faiss.IndexFlatL2(embedding_dimension)
        faiss_index.add(embeddings_matrix)
        logger.info("FAISS index updated.")

@api.route('/ask', methods=['POST'])
def ask_question():
    global documents, faiss_index
    question = request.json.get('question', '')
    logger.info(f"Received question: {question}")

    if not question:
        return jsonify({"error": "No question provided"}), 400

    if not documents:
        response = llm.invoke(question)
        return jsonify({"answer": str(response.content)}), 200

    question_embedding = embeddings.embed_query(question)
    distances, indices = faiss_index.search(np.array([question_embedding]), k=5)
    relevant_docs = [documents[i] for i in indices[0]]

    context = " ".join([doc.page_content for doc in relevant_docs])
    sources = [doc.metadata["source"] for doc in relevant_docs]

    prompt = (
        "Greet the user as Hello!\n"
        "You are a knowledgeable assistant. Use the context below to answer the question accurately. If you do not find any context, please use your own data.\n"
        "Context:\n"
        f"{context}\n\n"
        f"Question: {question}\n"
        "Answer:\n"
        "Sources: {', '.join(sources)}"
    )

    response = llm.invoke(prompt)
    logger.info("Question processed and answer generated.")
    return jsonify({"answer": str(response.content)}), 200

def process_text(text, source, chunk_size=1000, chunk_overlap=200):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_text(text)
    return [Document(page_content=chunk, metadata={"source": source}) for chunk in chunks]

def process_pdf(content, filename, chunk_size=1000, chunk_overlap=200):
    doc = fitz.open(stream=content, filetype="pdf")
    concatenated_content = ""
    for page in doc:
        concatenated_content += page.get_text()
    doc.close()
    return process_text(concatenated_content, filename, chunk_size, chunk_overlap)

if __name__ == '__main__':
    pass