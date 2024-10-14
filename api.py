## summary of the codebase below #
## The api.py file sets up a Flask application with three main endpoints:
# #1. /upload: Accepts file uploads (txt, docx, xlsx, pptx, pdf), processes the content, and stores it for later querying.
# #2. /process_urls: Accepts a list of URLs, extracts text from the web pages, and adds the content to the document store.
# #3 /ask: Takes a user question, retrieves relevant documents, and generates an answer using a language model, returning it as a JSON response.
#################

from flask import Flask, request, jsonify
import os
from io import BytesIO
import numpy as np
import faiss
from bs4 import BeautifulSoup
import requests
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import fitz  # PyMuPDF for PDF processing
from dotenv import load_dotenv
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
load_dotenv()

# Load environment variables
os.environ["AZURE_OPENAI_API_KEY"] = os.getenv("AZURE_OPENAI_API_KEY_GPT4")

llm = AzureChatOpenAI(
    openai_api_version=os.getenv("OPENAI_API_VERSION"),
    azure_deployment=os.getenv("AZURE_DEPLOYMENT_GPT4"),
    model_name=os.getenv("MODEL_NAME"),
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT_GPT4"),
    temperature=0.0,
)

os.environ["AZURE_OPENAI_API_KEY"] = os.getenv("AZURE_OPENAI_API_KEY_EMBEDDING")
os.environ["AZURE_OPENAI_ENDPOINT"] = os.getenv("AZURE_OPENAI_ENDPOINT_EMBEDDING")
os.environ["AZURE_OPENAI_API_VERSION"] = os.getenv("AZURE_OPENAI_API_VERSION_EMBEDDING")

embeddings = AzureOpenAIEmbeddings(
    azure_deployment=os.getenv("AZURE_EMBEDDING_NAME"),
    openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION_EMBEDDING")
)

documents = []
faiss_index = None


def process_text(text, source, chunk_size=1000, chunk_overlap=200):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = text_splitter.split_text(text)
    return [Document(page_content=chunk, metadata={"source": source}) for chunk in chunks]


def process_pdf(content, filename, chunk_size=1000, chunk_overlap=200):
    doc = fitz.open(stream=content, filetype="pdf")
    concatenated_content = ""
    for page in doc:
        concatenated_content += page.get_text()
    doc.close()
    return process_text(concatenated_content, filename, chunk_size, chunk_overlap)


@app.route('/upload', methods=['POST'])
def upload_file():
    global documents, faiss_index
    logger.info("Received file upload request.")

    if 'file' not in request.files:
        logger.error("No file part in the request.")
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        logger.error("No selected file.")
        return jsonify({"error": "No selected file"}), 400

    try:
        content = file.read()
        file_extension = os.path.splitext(file.filename)[1].lower()

        logger.info(f"Processing file: {file.filename} with extension: {file_extension}")

        if file_extension == '.txt':
            new_docs = process_text(content.decode('utf-8'), file.filename)
        elif file_extension == '.docx':
            doc = DocxDocument(BytesIO(content))
            text = "\n".join([para.text for para in doc.paragraphs if para.text])
            new_docs = process_text(text, file.filename)
        elif file_extension == '.xlsx':
            df = pd.read_excel(BytesIO(content))
            text = df.to_string(index=False)
            new_docs = process_text(text, file.filename)
        elif file_extension == '.pptx':
            prs = Presentation(BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            new_docs = process_text(text, file.filename)
        elif file_extension == '.pdf':
            new_docs = process_pdf(content, file.filename)
        else:
            logger.error(f"Unsupported file type: {file_extension}")
            return jsonify({"error": f"Unsupported file type: {file_extension}"}), 400

        documents.extend(new_docs)
        update_faiss_index()
        logger.info(f"Successfully processed {file.filename}.")
        return jsonify({"message": "File processed successfully"}), 200

    except Exception as e:
        logger.exception("An error occurred while processing the file.")
        return jsonify({"error": "An error occurred while processing the file."}), 500

@app.route('/process_urls', methods=['POST'])
def process_urls():
    global documents, faiss_index
    urls = request.json.get('urls', [])
    errors = []

    for url in urls:
        if not url.startswith(('http://', 'https://')):
            errors.append(f"Invalid URL '{url}': No scheme supplied. Perhaps you meant 'https://{url}'?")
            continue

        try:
            response = requests.get(url, verify=False)
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                text_content = soup.get_text()
                new_docs = process_text(text_content, url)
                documents.extend(new_docs)
            else:
                errors.append(f"Error fetching URL '{url}': Status code {response.status_code}.")
        except Exception as e:
            errors.append(f"Error processing URL '{url}': {str(e)}")

    update_faiss_index()

    if errors:
        return jsonify({"error": errors}), 400  # Return errors with 400 status code

    return jsonify({"message": "URLs processed successfully."}), 200


def update_faiss_index():
    global faiss_index
    if documents:
        embeddings_matrix = np.array(embeddings.embed_documents([doc.page_content for doc in documents]))
        embedding_dimension = embeddings_matrix.shape[1]
        faiss_index = faiss.IndexFlatL2(embedding_dimension)
        faiss_index.add(embeddings_matrix)
        logger.info("FAISS index updated.")


@app.route('/ask', methods=['POST'])
def ask_question():
    global documents, faiss_index
    question = request.json.get('question', '')

    logger.info(f"Received question: {question}")

    if not question:
        logger.error("No question provided.")
        return jsonify({"error": "No question provided"}), 400

    if not documents:
        logger.info("No documents uploaded, using LLM directly.")
        response = llm.invoke(question)
        return jsonify({"answer": str(response.content)}), 200

    question_embedding = embeddings.embed_query(question)
    distances, indices = faiss_index.search(np.array([question_embedding]), k=5)
    relevant_docs = [documents[i] for i in indices[0]]

    context = " ".join([doc.page_content for doc in relevant_docs])
    sources = [doc.metadata["source"] for doc in relevant_docs]

    prompt = (
        "Greet the user as Hello!\n" 
        "You are a knowledgeable assistant. Use the context below to answer the question accurately. If you do not find any context, please use your own data.\n"
        "Apply your intelligence on the context provided below \n"
        "Context:\n"
        f"{context}\n\n"
        f"Question: {question}\n"
        "Answer:\n"
        f"Sources: {', '.join(sources)}"
    )

    response = llm.invoke(prompt)
    logger.info("Question processed and answer generated.")
    return jsonify({"answer": str(response.content)}), 200


if __name__ == '__main__':
    app.run(debug=True)

